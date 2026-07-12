from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BG_DARK = RGBColor(0x0A, 0x0A, 0x0F)
BG_CARD = RGBColor(0x14, 0x14, 0x1E)
PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
BLUE = RGBColor(0x60, 0xA5, 0xFA)
GREEN = RGBColor(0x34, 0xD3, 0x99)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xA1, 0xA1, 0xAA)
DARK_GRAY = RGBColor(0x71, 0x71, 0x7A)
LIGHT_BG = RGBColor(0x1C, 0x1C, 0x28)
BORDER = RGBColor(0x2A, 0x2A, 0x38)

def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font_name='Inter'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf

def add_rect(slide, left, top, width, height, fill_color=None, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or BG_CARD
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_bullet_text(tf, text, level=0, font_size=14, color=GRAY, bold=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Inter'
    p.level = level
    return p

def add_circle(slide, left, top, size, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

# ========== SLIDE 1: TITLE ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_bg(slide)

# Badge
badge = add_rect(slide, Inches(5.2), Inches(1.5), Inches(3), Inches(0.45), PURPLE, PURPLE)
badge.fill.solid()
badge.fill.fore_color.rgb = RGBColor(0x1E, 0x10, 0x3A)
badge.line.color.rgb = RGBColor(0x4C, 0x1D, 0x95)
tf = badge.text_frame
tf.paragraphs[0].text = "🍏  v1 now in beta — Mac desktop"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = PURPLE
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.paragraphs[0].font.name = 'Inter'

# Title
add_text_box(slide, Inches(1.5), Inches(2.4), Inches(10.3), Inches(1.5),
    "NoteMeet", font_size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide, Inches(2.5), Inches(3.6), Inches(8.3), Inches(1.2),
    "AI Notepad for Indian Teams", font_size=36, bold=False, color=PURPLE, align=PP_ALIGN.CENTER)

# Tagline
add_text_box(slide, Inches(2), Inches(4.8), Inches(9.3), Inches(0.8),
    "Online meetings · WhatsApp calls · In-person conversations — one AI notepad.\nNo bots. No awkwardness. Just your notes, ready instantly.",
    font_size=18, bold=False, color=GRAY, align=PP_ALIGN.CENTER)

# Bottom tag
add_text_box(slide, Inches(3), Inches(6.5), Inches(7.3), Inches(0.5),
    "v1 in beta  ·  Mac only  ·  Free during beta  ·  Android soon",
    font_size=14, bold=False, color=DARK_GRAY, align=PP_ALIGN.CENTER)


# ========== SLIDE 2: PROBLEM ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(2), Inches(0.4),
    "THE PROBLEM", font_size=12, bold=True, color=PURPLE)
add_text_box(slide, Inches(0.8), Inches(1), Inches(11), Inches(0.8),
    "Indian professionals have no tool that actually works for them",
    font_size=36, bold=True, color=WHITE)

# Problem cards
problems = [
    ("🤖", "Meeting bots are invasive", "Otter, Fireflies join as participants. 68% of Indian companies ban external bots on calls."),
    ("💸", "Global tools cost $20/mo", "₹1,500-1,700/mo is 5-10x what Indian users will pay. Granola: $20. Otter: $17. Fireflies: $18."),
    ("🇮🇳", "No Indian language support", "Hinglish is the default in Indian meetings. Global tools' transcription accuracy drops 40% on code-switched speech."),
    ("📞", "WhatsApp calls leave no record", "India's primary business communication tool is WhatsApp. Zero tools capture WhatsApp call notes."),
    ("🔒", "Data privacy concerns", "DPDP Act 2023 requires data localization. Most meeting tools store data on US servers."),
    ("📝", "Manual notes are broken", "60% still take manual notes. 40% of action items are forgotten within 24 hours."),
]

for i, (icon, title, desc) in enumerate(problems):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(2.1 + row * 2.5)

    card = add_rect(slide, x, y, Inches(3.7), Inches(2.1), LIGHT_BG, BORDER)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.25), Inches(0.5), Inches(0.5),
        icon, font_size=28, color=WHITE)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.85), Inches(3.1), Inches(0.4),
        title, font_size=16, bold=True, color=WHITE)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.3), Inches(3.1), Inches(0.7),
        desc, font_size=12, color=GRAY)


# ========== SLIDE 3: SOLUTION ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(2), Inches(0.4),
    "THE SOLUTION", font_size=12, bold=True, color=PURPLE)
add_text_box(slide, Inches(0.8), Inches(1), Inches(11), Inches(0.8),
    "Meet NoteMeet", font_size=36, bold=True, color=WHITE)

# Three columns
for i, (title, desc, icon) in enumerate([
    ("Online Meetings", "Captures Zoom, Google Meet, Teams, Webex. Runs on your device — no bot joins your call. One click, perfect notes.", "🖥️"),
    ("In-Person & WhatsApp", "Open app, tap record. Works for client meetings, coffee chats, and WhatsApp/VoIP calls. Offline mode included.", "📱"),
    ("AI Brain", "Instant summaries, action items, decisions. Searchable history. Chat: 'What did I promise last week?'", "🧠"),
]):
    x = Inches(0.8 + i * 4.1)

    card = add_rect(slide, x, Inches(2.1), Inches(3.7), Inches(3.0), LIGHT_BG, BORDER)
    add_text_box(slide, x + Inches(1.4), Inches(2.4), Inches(1), Inches(0.8),
        icon, font_size=42, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), Inches(3.2), Inches(3.1), Inches(0.4),
        title, font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), Inches(3.7), Inches(3.1), Inches(1.2),
        desc, font_size=13, color=GRAY, align=PP_ALIGN.CENTER)

# Bottom highlight
highlight = add_rect(slide, Inches(1.5), Inches(5.8), Inches(10.3), Inches(1.0), RGBColor(0x1E, 0x10, 0x3A), RGBColor(0x4C, 0x1D, 0x95))
add_text_box(slide, Inches(1.8), Inches(5.95), Inches(9.7), Inches(0.7),
    "Key difference: No bot joins your meetings. NoteMeet listens from your device — just like you do.",
    font_size=16, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)


# ========== SLIDE 4: WHY INDIA ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(2), Inches(0.4),
    "WHY INDIA", font_size=12, bold=True, color=PURPLE)
add_text_box(slide, Inches(0.8), Inches(1), Inches(11), Inches(0.8),
    "Built from the ground up for Indian teams",
    font_size=36, bold=True, color=WHITE)

items = [
    ("🇮🇳", "Hinglish is first-class", "Understands Hindi, Tamil, Telugu, Bengali, and mixed Hinglish. No other tool handles code-switching."),
    ("📞", "WhatsApp-native", "Primary business communication in India is WhatsApp. NoteMeet captures WhatsApp calls automatically."),
    ("🔒", "DPDP Act Compliant", "Data stays on your device or on Indian servers. Fully compliant with India's data protection laws."),
    ("💰", "7x cheaper", "₹199/mo vs ₹1,500-1,700 for Fireflies/Otter/Granola. Pricing built for the Indian market."),
    ("📱", "Mobile-first, offline-first", "800M+ smartphone users in India. Works offline. Optimized for Jio/Airtel networks."),
    ("🏢", "Startup & GCC focused", "Built for the 45,000+ funded startup execs and 300K+ tech managers who live in meetings."),
]

for i, (icon, title, desc) in enumerate(items):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(2.1 + row * 2.5)

    card = add_rect(slide, x, y, Inches(3.7), Inches(2.1), LIGHT_BG, BORDER)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.25), Inches(0.5), Inches(0.5),
        icon, font_size=28, color=WHITE)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.85), Inches(3.1), Inches(0.4),
        title, font_size=16, bold=True, color=WHITE)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.3), Inches(3.1), Inches(0.7),
        desc, font_size=12, color=GRAY)


# ========== SLIDE 5: MARKET ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(2), Inches(0.4),
    "MARKET OPPORTUNITY", font_size=12, bold=True, color=PURPLE)
add_text_box(slide, Inches(0.8), Inches(1), Inches(11), Inches(0.8),
    "Large, underserved, and ready for an India-first solution",
    font_size=36, bold=True, color=WHITE)

# TAM boxes
tam_data = [
    ("₹24,000 Cr", "TAM", "100M knowledge workers in India"),
    ("₹7,200 Cr", "SAM", "30M professionals, 5+ meetings/week"),
    ("₹480 Cr", "SOM", "2M early adopters (tech/startups/consulting)"),
]

for i, (val, label, desc) in enumerate(tam_data):
    x = Inches(0.8 + i * 4.1)
    box = add_rect(slide, x, Inches(2.1), Inches(3.7), Inches(1.6), LIGHT_BG, BORDER)
    add_text_box(slide, x + Inches(0.3), Inches(2.3), Inches(3.1), Inches(0.7),
        val, font_size=36, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), Inches(2.85), Inches(3.1), Inches(0.3),
        label, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), Inches(3.15), Inches(3.1), Inches(0.4),
        desc, font_size=12, color=GRAY, align=PP_ALIGN.CENTER)

# Target segment
segment = add_rect(slide, Inches(0.8), Inches(4.2), Inches(11.7), Inches(2.8), LIGHT_BG, BORDER)
add_text_box(slide, Inches(1.2), Inches(4.4), Inches(11), Inches(0.4),
    "Primary Beachhead: Startup Founders & Tech Managers", font_size=20, bold=True, color=WHITE)

segments = [
    ("45,000", "Founder/C-suite at\nfunded startups"),
    ("2,00,000", "Tech Managers\n(EM, PM, Design Lead)"),
    ("1,50,000", "Tech Managers\nat GCCs (1,650+ in India)"),
    ("2,00,000", "Senior ICs at\ngrowth companies"),
]

for i, (num, label) in enumerate(segments):
    x = Inches(1.2 + i * 2.9)
    add_text_box(slide, x, Inches(5.0), Inches(2.5), Inches(0.5),
        num, font_size=28, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(5.55), Inches(2.5), Inches(0.6),
        label, font_size=11, color=GRAY, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.2), Inches(6.4), Inches(11), Inches(0.3),
    "Need just 0.17% penetration to reach 1,000 paid users → ₹2L MRR",
    font_size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)


# ========== SLIDE 6: COMPETITION ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(2), Inches(0.4),
    "COMPETITIVE LANDSCAPE", font_size=12, bold=True, color=PURPLE)
add_text_box(slide, Inches(0.8), Inches(1), Inches(11), Inches(0.8),
    "Everyone is either expensive, bot-based, or ignores India",
    font_size=36, bold=True, color=WHITE)

# Table-like comparison
headers = ["Product", "Price (₹/mo)", "Bot?", "Indian Langs", "WhatsApp", "Mobile"]
rows = [
    ["Otter.ai", "₹1,400", "Yes ❌", "No ❌", "No ❌", "Yes"],
    ["Fireflies.ai", "₹1,500", "Yes ❌", "Limited", "No ❌", "Yes"],
    ["Granola", "₹1,650", "No ✅", "No ❌", "No ❌", "iOS only"],
    ["Talat", "Free", "No ✅", "Hindi only", "No ❌", "No ❌"],
    ["Sutra AI", "₹399/60min", "No ✅", "10 langs ✅", "No ❌", "No ❌"],
    ["NoteMeet", "₹199", "No ✅", "10+ langs ✅", "Yes ✅", "Yes ✅"],
]

col_widths = [Inches(2.3), Inches(1.5), Inches(1.3), Inches(1.8), Inches(1.5), Inches(1.3)]
table_top = Inches(2.2)
table_left = Inches(1.2)

for j, header in enumerate(headers):
    x = table_left + sum(w for w in col_widths[:j])
    box = add_rect(slide, x, table_top, col_widths[j], Inches(0.5), RGBColor(0x2A, 0x2A, 0x38), None)
    add_text_box(slide, x, table_top + Inches(0.08), col_widths[j], Inches(0.35),
        header, font_size=12, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

for i, row in enumerate(rows):
    y = table_top + Inches(0.5 + i * 0.5)
    bg = RGBColor(0x1E, 0x10, 0x3A) if i == len(rows) - 1 else LIGHT_BG
    border = RGBColor(0x4C, 0x1D, 0x95) if i == len(rows) - 1 else None
    for j, cell in enumerate(row):
        x = table_left + sum(w for w in col_widths[:j])
        box = add_rect(slide, x, y, col_widths[j], Inches(0.45), bg, border)
        color = GREEN if (i == len(rows) - 1 and cell == "Yes ✅") else WHITE if i == len(rows) - 1 else GRAY
        bold = True if i == len(rows) - 1 else False
        add_text_box(slide, x, y + Inches(0.07), col_widths[j], Inches(0.35),
            cell, font_size=11, bold=bold, color=color, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.2), Inches(5.8), Inches(11), Inches(0.8),
    "NoteMeet is the ONLY product combining: no-bot + WhatsApp capture + Indian languages + mobile app + ₹199 pricing",
    font_size=15, bold=True, color=GREEN, align=PP_ALIGN.CENTER)


# ========== SLIDE 7: WHY NO PHONE CALLS ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(2), Inches(0.4),
    "THE ELEPHANT IN THE ROOM", font_size=12, bold=True, color=PURPLE)
add_text_box(slide, Inches(0.8), Inches(1), Inches(11), Inches(0.8),
    "Why doesn't anyone support mobile call recording?",
    font_size=36, bold=True, color=WHITE)
add_text_box(slide, Inches(0.8), Inches(1.7), Inches(11), Inches(0.5),
    "Apple and Google have made it technically impossible for third-party apps. Here's exactly why.",
    font_size=16, color=GRAY)

# iOS column
ios_card = add_rect(slide, Inches(0.8), Inches(2.4), Inches(5.7), Inches(4.6), LIGHT_BG, BORDER)
add_text_box(slide, Inches(1.1), Inches(2.6), Inches(5.1), Inches(0.4),
    "🍎  iPhone (iOS)", font_size=20, bold=True, color=WHITE)

ios_reasons = [
    ("🔒  iOS blocks ALL call audio access",
     "Apple provides zero APIs to capture cellular call audio. It's a deliberate privacy decision — no third-party app can listen to your phone calls."),
    ("🎤  Mic is disabled during calls",
     "When a cellular call is active, iOS stops giving microphone access to other apps. Any app that tries to record gets silence."),
    ("📁  iOS 18 recording is sandboxed",
     "Apple's built-in call recording (iOS 18.1+) exists, but recordings are locked inside the Phone app. No share button. No export. No API to access them."),
    ("📞  VoIP calls (WhatsApp) work — partially",
     "WhatsApp calls on iOS go through VoIP APIs. ReplayKit can capture screen+audio but shows a prominent green 'Recording' banner — awkward in a business call."),
]

tf = add_text_box(slide, Inches(1.1), Inches(3.15), Inches(5.1), Inches(3.6), "", font_size=13, color=GRAY)
for title, desc in ios_reasons:
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Inter'
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(11)
    p2.font.color.rgb = GRAY
    p2.font.name = 'Inter'
    p2.space_after = Pt(10)

# Android column
android_card = add_rect(slide, Inches(6.8), Inches(2.4), Inches(5.7), Inches(4.6), LIGHT_BG, BORDER)
add_text_box(slide, Inches(7.1), Inches(2.6), Inches(5.1), Inches(0.4),
    "🤖  Android", font_size=20, bold=True, color=WHITE)

android_reasons = [
    ("🔇  VOICE_CALL source blocked in Android 10+",
     "Android used to have a 'VOICE_CALL' audio source for recording calls. Google removed it in Android 10. Now it returns encrypted/empty audio."),
    ("🚫  Play Store banned call recording apps",
     "Since 2023, Google Play blocks apps that use AccessibilityService for call recording. You can only sideload — impossible for mass market."),
    ("🏭  OEMs have their own recorders",
     "Xiaomi, Samsung, OnePlus ship built-in call recording in India ROMs. But it's locked to their dialer — third-party apps can't access it."),
    ("🛠️  The only way? Root the phone",
     "Apps like ACR need root access. Less than 0.5% of Indian users root their phones. Not a viable product strategy."),
]

tf = add_text_box(slide, Inches(7.1), Inches(3.15), Inches(5.1), Inches(3.6), "", font_size=13, color=GRAY)
for title, desc in android_reasons:
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Inter'
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(11)
    p2.font.color.rgb = GRAY
    p2.font.name = 'Inter'
    p2.space_after = Pt(10)

# Bottom conclusion
conclusion = add_rect(slide, Inches(0.8), Inches(7.25), Inches(11.7), Inches(0.5), RGBColor(0x1E, 0x10, 0x3A), RGBColor(0x4C, 0x1D, 0x95))
add_text_box(slide, Inches(1.1), Inches(7.3), Inches(11.1), Inches(0.35),
    "Granola tries but fails: iOS phone calls India is blocked. WhatsApp mobile calls can't be captured cleanly. The industry's best answer is 'desktop-only.'",
    font_size=12, bold=False, color=PURPLE, align=PP_ALIGN.CENTER)


# ========== SLIDE 8: OUR APPROACH TO WHATSAPP ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(2), Inches(0.4),
    "OUR APPROACH", font_size=12, bold=True, color=PURPLE)
add_text_box(slide, Inches(0.8), Inches(1), Inches(11), Inches(0.8),
    "How we handle WhatsApp & phone calls (honestly)",
    font_size=36, bold=True, color=WHITE)

# What we DO support
do_card = add_rect(slide, Inches(0.8), Inches(2.1), Inches(5.7), Inches(3.0), LIGHT_BG, BORDER)
add_text_box(slide, Inches(1.1), Inches(2.3), Inches(5.1), Inches(0.4),
    "✅  What We Support", font_size=18, bold=True, color=GREEN)

do_items = [
    ("Desktop WhatsApp Web calls", "System audio capture on Mac/Windows — clean, no banner, works automatically. This covers most work-related WhatsApp calls in India."),
    ("Online meetings", "Zoom, Google Meet, Teams, Webex — via system audio. No bot joins your call."),
    ("In-person conversations", "Mobile app mic recording. Works offline. Perfect for client meetings."),
    ("Outbound calls via in-app dialer (future)", "SIP-based dialer (like Truecaller) — user dials through our app, we merge recording. Works on iOS + Android."),
]

tf = add_text_box(slide, Inches(1.1), Inches(2.9), Inches(5.1), Inches(2.0), "", font_size=12, color=GRAY)
for title, desc in do_items:
    p = tf.add_paragraph()
    p.text = "→  " + title
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Inter'
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(11)
    p2.font.color.rgb = GRAY
    p2.font.name = 'Inter'
    p2.space_after = Pt(8)

# What we DON'T support
dont_card = add_rect(slide, Inches(6.8), Inches(2.1), Inches(5.7), Inches(3.0), LIGHT_BG, BORDER)
add_text_box(slide, Inches(7.1), Inches(2.3), Inches(5.1), Inches(0.4),
    "❌  What We Cannot Do (honest)", font_size=18, bold=True, color=RGBColor(0xEF, 0x44, 0x44))

dont_items = [
    ("Mobile WhatsApp calls (no workaround)", "iOS blocks it. Android blocks it. Even Granola, Otter, Fireflies can't do it. We won't pretend otherwise."),
    ("Incoming native phone calls", "iOS: no API. Android: blocked. Granola tried via SIP dialer but India is banned. No third-party app in the world supports this cleanly."),
    ("Recording without user awareness", "All recording requires user intent (tap record). No always-listening. Good privacy = deliberate action."),
]

tf = add_text_box(slide, Inches(7.1), Inches(2.9), Inches(5.1), Inches(2.0), "", font_size=12, color=GRAY)
for title, desc in dont_items:
    p = tf.add_paragraph()
    p.text = "✗  " + title
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xEF, 0x44, 0x44)
    p.font.name = 'Inter'
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(11)
    p2.font.color.rgb = GRAY
    p2.font.name = 'Inter'
    p2.space_after = Pt(8)

# Market context
context = add_rect(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.8), LIGHT_BG, BORDER)
add_text_box(slide, Inches(1.1), Inches(5.7), Inches(11.1), Inches(0.4),
    "📊  Market Context — Why This Still Wins", font_size=18, bold=True, color=BLUE)

context_items = [
    "🗣️  90%+ of professional conversations in India happen on: Zoom/Meet (desktop) + in-person + WhatsApp Web (desktop) — all three we support.",
    "📱  Mobile WhatsApp calls are mostly personal (family, friends). Professionals use WhatsApp Web on laptop for work calls.",
    "🏆  Granola has a $125M valuation and ALSO cannot do mobile WhatsApp calls. Nobody can. This is not a weakness — it's the industry reality.",
    "💡  When Apple/Google eventually open call recording APIs (pressure from regulators), we'll be first to integrate. Until then, we ship what works.",
]

tf = add_text_box(slide, Inches(1.1), Inches(6.2), Inches(11.1), Inches(1.0), "", font_size=12, color=GRAY)
for item in context_items:
    p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY
    p.font.name = 'Inter'
    p.space_after = Pt(4)


# ========== SLIDE 9: BUSINESS MODEL ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(2), Inches(0.4),
    "BUSINESS MODEL", font_size=12, bold=True, color=PURPLE)
add_text_box(slide, Inches(0.8), Inches(1), Inches(11), Inches(0.8),
    "Freemium SaaS — radically affordable for India",
    font_size=36, bold=True, color=WHITE)

# Pricing cards
pricing = [
    ("Free", "₹0/mo", "10 meetings\n7-day history\nBasic summaries\nSingle device"),
    ("Pro", "₹199/mo", "Unlimited meetings\n1-year history\nAI Chat + WhatsApp\nCalendar sync + Briefs\n10 Indian languages"),
    ("Business", "₹599/user/mo", "Everything in Pro\nTeam dashboard\nAdmin controls & SSO\nCRM integrations\nPriority support"),
]

for i, (plan, price, features) in enumerate(pricing):
    x = Inches(0.8 + i * 4.1)
    is_featured = i == 1

    card = add_rect(slide, x, Inches(2.1), Inches(3.7), Inches(3.8),
        RGBColor(0x1E, 0x10, 0x3A) if is_featured else LIGHT_BG,
        RGBColor(0x4C, 0x1D, 0x95) if is_featured else BORDER)

    if is_featured:
        tag = add_rect(slide, x + Inches(0.8), Inches(1.85), Inches(2.1), Inches(0.35),
            RGBColor(0x4C, 0x1D, 0x95), None)
        tag.fill.solid()
        tag.fill.fore_color.rgb = RGBColor(0x4C, 0x1D, 0x95)
        add_text_box(slide, x + Inches(0.8), Inches(1.88), Inches(2.1), Inches(0.3),
            "★ Most Popular", font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide, x + Inches(0.3), Inches(2.3 + (0.25 if is_featured else 0)),
        Inches(3.1), Inches(0.3), plan, font_size=14, bold=True, color=DARK_GRAY)
    add_text_box(slide, x + Inches(0.3), Inches(2.7 + (0.25 if is_featured else 0)),
        Inches(3.1), Inches(0.5), price, font_size=28, bold=True, color=WHITE)

    # Features
    feat_tf = add_text_box(slide, x + Inches(0.3), Inches(3.4 + (0.25 if is_featured else 0)),
        Inches(3.1), Inches(2.0), "", font_size=12, color=GRAY)
    for line in features.split('\n'):
        add_bullet_text(feat_tf, "✓  " + line, font_size=12, color=GRAY)

# Revenue
rev_card = add_rect(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.8), LIGHT_BG, BORDER)
add_text_box(slide, Inches(1.2), Inches(6.35), Inches(11), Inches(0.5),
    "Revenue at scale:  2,000 users → ₹6L MRR  |  10,000 users → ₹32L MRR  |  50,000 users → ₹1.6Cr MRR",
    font_size=15, bold=True, color=GREEN, align=PP_ALIGN.CENTER)


# ========== SLIDE 8: BUILD PLAN ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(2), Inches(0.4),
    "BUILD PLAN", font_size=12, bold=True, color=PURPLE)
add_text_box(slide, Inches(0.8), Inches(1), Inches(11), Inches(0.8),
    "v1 is shipping. Here's what's next.",
    font_size=36, bold=True, color=WHITE)

# Timeline — with v1 completed
phases = [
    ("✅ COMPLETED", "v1 — Mac Desktop", "System audio capture\nWhisper on-device transcription\nLLM note generation\nBasic UI + meeting history"),
    ("⚠️ IN PROGRESS", "Calendar + Chat", "Google Calendar sync\nMeeting detection\nAI chat interface\nUser auth & sync"),
    ("📅 NEXT", "Mobile Apps", "iOS mic recording\nAndroid capture\nIn-person meeting mode\nCross-device sync"),
    ("📅 NEXT", "WhatsApp + Scale", "WhatsApp desktop capture\nOnboarding & referral\nAndroid + Windows\n₹199/mo launch"),
]

for i, (phase, title, desc) in enumerate(phases):
    x = Inches(0.8 + i * 3.15)
    card = add_rect(slide, x, Inches(2.1), Inches(2.85), Inches(3.0), LIGHT_BG, BORDER)
    add_text_box(slide, x + Inches(0.3), Inches(2.25), Inches(2.25), Inches(0.3),
        phase, font_size=11, bold=True, color=GREEN if 'COMPLETED' in phase else PURPLE)
    add_text_box(slide, x + Inches(0.3), Inches(2.6), Inches(2.25), Inches(0.4),
        title, font_size=18, bold=True, color=WHITE)
    add_text_box(slide, x + Inches(0.3), Inches(3.1), Inches(2.25), Inches(1.2),
        desc, font_size=12, color=GRAY)

# Cost box
cost = add_rect(slide, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.5), LIGHT_BG, BORDER)
add_text_box(slide, Inches(1.2), Inches(5.75), Inches(11), Inches(0.3),
    "Total Cash Outlay to MVP: ₹59,000 ($710)", font_size=20, bold=True, color=WHITE)

cost_items = [
    ("₹7,000", "Apple Developer\nAccount (annual)"),
    ("₹2,000", "Google Play\n(one-time)"),
    ("₹30,000", "API credits\n(dev & testing)"),
    ("₹20,000", "Cloud infra\n(first 3 months)"),
]
for i, (amount, label) in enumerate(cost_items):
    x = Inches(1.2 + i * 2.9)
    add_text_box(slide, x, Inches(6.2), Inches(2.5), Inches(0.4),
        amount, font_size=22, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(6.6), Inches(2.5), Inches(0.4),
        label, font_size=10, color=GRAY, align=PP_ALIGN.CENTER)


# ========== SLIDE 9: ASK ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(2), Inches(0.4),
    "INVESTMENT ASK", font_size=12, bold=True, color=PURPLE)
add_text_box(slide, Inches(0.8), Inches(1), Inches(11), Inches(0.8),
    "Let's build India's AI notepad",
    font_size=36, bold=True, color=WHITE)

# Three stages
stages = [
    ("Pre-Seed", "₹15-25L\n($18-30k)", [
        "API credits for MVP",
        "Cloud infrastructure",
        "Marketing & launch",
        "Legal & compliance",
    ]),
    ("Seed", "₹1-2 Cr\n($120-240k)", [
        "Hire 2 devs + 1 designer",
        "Full-time founder salary",
        "India-first LLM training",
        "Enterprise pilots",
    ]),
    ("Series A", "₹5-10 Cr\n($600k-1.2M)", [
        "Scale to 500K users",
        "On-device AI team",
        "Enterprise sales team",
        "Data center in India",
    ]),
]

for i, (round_name, amount, uses) in enumerate(stages):
    x = Inches(0.8 + i * 4.1)
    card = add_rect(slide, x, Inches(2.1), Inches(3.7), Inches(4.5), LIGHT_BG, BORDER)
    add_text_box(slide, x + Inches(0.3), Inches(2.3), Inches(3.1), Inches(0.3),
        round_name, font_size=14, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), Inches(2.7), Inches(3.1), Inches(0.8),
        amount, font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    tf = add_text_box(slide, x + Inches(0.3), Inches(3.7), Inches(3.1), Inches(2.5),
        "", font_size=13, color=GRAY)
    for use in uses:
        add_bullet_text(tf, "→  " + use, font_size=13, color=GRAY)

# Milestones
milestones = add_rect(slide, Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.5), LIGHT_BG, BORDER)
add_text_box(slide, Inches(1.2), Inches(6.95), Inches(11), Inches(0.35),
    "Milestones:  1K paid users @ ₹2L MRR → Pre-seed   |   5K users @ ₹10L MRR → Seed   |   25K users @ ₹50L MRR → Series A",
    font_size=13, bold=False, color=BLUE, align=PP_ALIGN.CENTER)


# ========== SLIDE 10: WHY NOW + CLOSING ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

# Decorative orbs
add_circle(slide, Inches(-2), Inches(-2), Inches(8), RGBColor(0x2E, 0x10, 0x3A))
add_circle(slide, Inches(9), Inches(4), Inches(6), RGBColor(0x10, 0x2E, 0x3A))

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(2), Inches(0.4),
    "WHY NOW", font_size=12, bold=True, color=PURPLE)
add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.0),
    "The window is open — but closing fast",
    font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

reasons = [
    "✅ DPDP Act enforcement (2025) makes US-hosted tools risky for Indian enterprises",
    "✅ WhatsApp Business has 200M+ Indian users — voice notes need a capture tool",
    "✅ On-device AI (Whisper, Llama 3) now runs on phones — privacy-first is finally feasible",
    "✅ Competitors (Talat, Sutra) launched recently — no one has won yet",
    "✅ $125M into Granola validates the category — apply the model to India at 1/10th the price",
]

tf = add_text_box(slide, Inches(2), Inches(2.8), Inches(9.3), Inches(3.5), "", font_size=16, color=GRAY)
for r in reasons:
    p = tf.add_paragraph()
    p.text = r
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY
    p.font.name = 'Inter'
    p.space_after = Pt(12)

# Bottom CTA
cta = add_rect(slide, Inches(3), Inches(6.2), Inches(7.3), Inches(0.8), RGBColor(0x1E, 0x10, 0x3A), RGBColor(0x4C, 0x1D, 0x95))
add_text_box(slide, Inches(3.3), Inches(6.35), Inches(6.7), Inches(0.5),
    "Let's build the AI notepad that India deserves.", font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(3.5), Inches(7.0), Inches(6.3), Inches(0.4),
    "Solo founder  ·  8-week MVP  ·  ₹59k build cost  ·  ₹199/mo pricing",
    font_size=13, color=DARK_GRAY, align=PP_ALIGN.CENTER)


# Save
output_path = os.path.expanduser("~/Desktop/NoteMeet_Investor_Pitch.pptx")
prs.save(output_path)
print(f"✅ Saved to {output_path}")
print(f"   Slides: {len(prs.slides)}")
