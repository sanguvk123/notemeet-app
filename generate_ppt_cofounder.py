from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(0x00, 0x00, 0x00)
SURFACE = RGBColor(0x0A, 0x0A, 0x0A)
SURFACE2 = RGBColor(0x14, 0x14, 0x14)
SURFACE3 = RGBColor(0x1C, 0x1C, 0x1C)
BORDER = RGBColor(0x1F, 0x1F, 0x1F)
BORDER2 = RGBColor(0x2A, 0x2A, 0x2A)
TEXT = RGBColor(0xF5, 0xF5, 0xF5)
TEXT2 = RGBColor(0xA0, 0xA0, 0xA0)
TEXT3 = RGBColor(0x60, 0x60, 0x60)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xE0, 0x40, 0x40)

def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_tb(slide, l, t, w, h, text, size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Inter'
    p.alignment = align
    return tf

def add_card(slide, l, t, w, h, fill=SURFACE2, border=BORDER):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def add_bullet(tf, text, size=13, color=TEXT2, bold=False, indent=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Inter'
    p.level = indent
    return p

def section_header(slide, label, title):
    add_tb(slide, Inches(0.8), Inches(0.4), Inches(3), Inches(0.3),
        label, size=10, bold=True, color=TEXT3, align=PP_ALIGN.LEFT)
    add_tb(slide, Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.6),
        title, size=28, bold=True, color=TEXT)

# ========== SLIDE 1: TITLE ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_tb(s, Inches(3), Inches(0.8), Inches(7.3), Inches(0.3),
    "NOTE MEET — FOR CO-FOUNDER EYES ONLY", size=10, bold=True, color=TEXT3, align=PP_ALIGN.CENTER)
add_tb(s, Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.2),
    "NoteMeet", size=64, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
add_tb(s, Inches(2), Inches(3.3), Inches(9.3), Inches(0.6),
    "AI Notepad for Indian Teams", size=28, bold=False, color=TEXT2, align=PP_ALIGN.CENTER)
add_tb(s, Inches(2), Inches(4.3), Inches(9.3), Inches(1.0),
    "A co-founder deep-dive: product, market, revenue, team, risks.",
    size=16, color=TEXT3, align=PP_ALIGN.CENTER)
add_tb(s, Inches(3.5), Inches(6.0), Inches(6.3), Inches(0.5),
    "Confidential  ·  July 2026  ·  v1 in beta  ·  Mac desktop",
    size=12, color=TEXT3, align=PP_ALIGN.CENTER)

# ========== SLIDE 2: WHAT IS NOTE MEET ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
section_header(s, "PRODUCT", "What is NoteMeet?")

add_tb(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.5),
    "A desktop app that listens to your meetings and writes notes for you. No bots. No setup. Just tap record.",
    size=15, color=TEXT2)

cards = [
    ("🎤", "System Audio Capture", "Captures Zoom, Meet, Teams, Webex, WhatsApp Web calls. Runs entirely on your device — no bot joins the meeting."),
    ("🧠", "On-Device AI (Whisper + LLM)", "Transcribes locally using OpenAI Whisper (ggml). Generates summaries, action items, decisions, and transcripts via LLM (OpenAI/OpenRouter)."),
    ("📅", "Google Calendar Sync", "Reads your calendar to detect meetings. Syncs events. Google OAuth with PKCE — tokens stored locally."),
    ("💬", "AI Chat Across All Notes", "Ask questions across every meeting: 'What did I promise last week?' or 'Summarize all action items.'"),
    ("🔒", "Privacy-First by Design", "Audio & transcripts stay on your machine. No cloud storage. DPDP Act compliant. No data leaves without your intent."),
    ("📱", "Mini Floating Recorder", "Record from anywhere with a compact overlay window. Start/stop without opening the main app."),
]

for i, (icon, title, desc) in enumerate(cards):
    col, row = i % 3, i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(2.4 + row * 2.5)
    add_card(s, x, y, Inches(3.7), Inches(2.1))
    add_tb(s, x + Inches(0.25), y + Inches(0.2), Inches(0.4), Inches(0.4), icon, size=22)
    add_tb(s, x + Inches(0.25), y + Inches(0.7), Inches(3.2), Inches(0.3), title, size=14, bold=True, color=TEXT)
    add_tb(s, x + Inches(0.25), y + Inches(1.1), Inches(3.2), Inches(0.8), desc, size=11, color=TEXT2)

# ========== SLIDE 3: CURRENT STATE ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
section_header(s, "STATUS", "Where we are right now (July 2026)")

add_card(s, Inches(0.8), Inches(1.7), Inches(5.7), Inches(5.2))
add_tb(s, Inches(1.1), Inches(1.9), Inches(5.1), Inches(0.3), "✅  What Works", size=16, bold=True, color=TEXT)

items_done = [
    "Mac desktop app (Tauri + React + Rust)",
    "System audio capture via cpal (CoreAudio)",
    "On-device Whisper transcription (ggml-base)",
    "LLM note generation (summary, action items, decisions, tone, speakers)",
    "AI Chat across all notes",
    "Google OAuth (PKCE + system browser + local TCP server)",
    "Google Calendar read & sync",
    "Google Calendar event creation",
    "Mini floating recorder overlay window",
    "Guest mode (full app without sign-in)",
    "Pure monochromatic black UI design",
    "52 unit tests (32 Rust + 20 frontend)",
]
tf = add_tb(s, Inches(1.1), Inches(2.4), Inches(5.1), Inches(4.2), "", size=12, color=TEXT2)
for item in items_done:
    add_bullet(tf, f"✓  {item}", size=11, color=TEXT2)

add_card(s, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.2))
add_tb(s, Inches(7.1), Inches(1.9), Inches(5.1), Inches(0.3), "❌  What's Missing", size=16, bold=True, color=RED)

items_missing = [
    "Windows & Linux builds",
    "Mobile apps (iOS & Android)",
    "Mobile WhatsApp call capture (OS-blocked — see slide 8)",
    "Cloud sync across devices",
    "User accounts & team management",
    "Payment / billing integration (Razorpay/Stripe)",
    "Onboarding flow & tutorial",
    "Email/password auth (Google-only currently)",
    "Offline-first reliability improvements",
    "Push-to-talk / always-listening mode",
    "Enterprise SSO & admin controls",
    "App Store distribution (needs Apple Developer Program)",
]
tf = add_tb(s, Inches(7.1), Inches(2.4), Inches(5.1), Inches(4.2), "", size=12, color=TEXT2)
for item in items_missing:
    add_bullet(tf, f"✗  {item}", size=11, color=RED)

# ========== SLIDE 4: PRICING ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
section_header(s, "BUSINESS MODEL", "Pricing Strategy")

pricing_tiers = [
    ("Free", "₹0/mo", [
        "10 meetings/month",
        "7-day history",
        "Basic summaries only",
        "Single device",
        "No calendar sync",
    ], False),
    ("Pro", "₹199/mo", [
        "Unlimited meetings",
        "1-year history",
        "AI Chat across all notes",
        "WhatsApp Web capture",
        "Calendar sync + Meeting detection",
        "10 Indian languages",
    ], True),
    ("Business", "₹599/user/mo", [
        "Everything in Pro",
        "Team dashboard",
        "Admin controls + SSO",
        "CRM integrations",
        "Priority support",
        "On-premise option",
    ], False),
]

for i, (plan, price, features, featured) in enumerate(pricing_tiers):
    x = Inches(0.8 + i * 4.1)
    fill = SURFACE3 if featured else SURFACE2
    border = BORDER2 if featured else BORDER
    add_card(s, x, Inches(1.6), Inches(3.7), Inches(4.8), fill, border)
    if featured:
        tag = add_card(s, x + Inches(0.8), Inches(1.35), Inches(2.1), Inches(0.35), SURFACE3, BORDER2)
        add_tb(s, x + Inches(0.8), Inches(1.38), Inches(2.1), Inches(0.3),
            "★ Recommended", size=10, bold=True, color=TEXT2, align=PP_ALIGN.CENTER)
    add_tb(s, x + Inches(0.3), Inches(1.8 + (0.2 if featured else 0)),
        Inches(3.1), Inches(0.3), plan, size=13, bold=True, color=TEXT3, align=PP_ALIGN.CENTER)
    add_tb(s, x + Inches(0.3), Inches(2.2 + (0.2 if featured else 0)),
        Inches(3.1), Inches(0.5), price, size=30, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    tf = add_tb(s, x + Inches(0.3), Inches(2.9 + (0.2 if featured else 0)),
        Inches(3.1), Inches(3.0), "", size=12, color=TEXT2)
    for f in features:
        add_bullet(tf, f"→  {f}", size=11, color=TEXT2)

add_card(s, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.6))
add_tb(s, Inches(1.2), Inches(6.68), Inches(11), Inches(0.35),
    "Why ₹199?  Comparable to 1 chai + samosa for two.  Otter charges ₹1,400.  Granola charges ₹1,650.  We're 7-8x cheaper.",
    size=13, bold=False, color=TEXT2, align=PP_ALIGN.CENTER)

# ========== SLIDE 5: TAM & REVENUE ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
section_header(s, "MARKET", "TAM, SAM, SOM & Revenue Projection")

# TAM boxes
tams = [
    ("₹24,000 Cr", "TAM", "100M knowledge workers in India\neven at ₹199/mo"),
    ("₹7,200 Cr", "SAM", "30M professionals\n5+ meetings/week"),
    ("₹480 Cr", "SOM", "2M early adopters\ntech/startups/consulting"),
]
for i, (val, label, desc) in enumerate(tams):
    x = Inches(0.8 + i * 4.1)
    add_card(s, x, Inches(1.5), Inches(3.7), Inches(1.5))
    add_tb(s, x + Inches(0.2), Inches(1.65), Inches(3.3), Inches(0.5),
        val, size=30, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    add_tb(s, x + Inches(0.2), Inches(2.15), Inches(3.3), Inches(0.25),
        label, size=12, bold=True, color=TEXT3, align=PP_ALIGN.CENTER)
    add_tb(s, x + Inches(0.2), Inches(2.4), Inches(3.3), Inches(0.4),
        desc, size=10, color=TEXT3, align=PP_ALIGN.CENTER)

# Revenue projection table
add_card(s, Inches(0.8), Inches(3.3), Inches(11.7), Inches(3.8))
add_tb(s, Inches(1.1), Inches(3.5), Inches(5), Inches(0.3),
    "Revenue Projection (Conservative)", size=14, bold=True, color=TEXT)

headers = ["Metric", "Year 1", "Year 2", "Year 3"]
rows = [
    ["Paid Users", "1,000", "10,000", "50,000"],
    ["Pro/Free Split", "60% / 40%", "50% / 50%", "40% / 60%"],
    ["Avg Rev/User (Pro)", "₹199/mo", "₹199/mo", "₹249/mo"],
    ["Blended ARPU", "₹119/mo", "₹100/mo", "₹100/mo"],
    ["Monthly Revenue", "₹1.2L ($1,440)", "₹10L ($12k)", "₹50L ($60k)"],
    ["Annual Revenue", "₹14.4L ($17k)", "₹1.2Cr ($144k)", "₹6Cr ($720k)"],
    ["Gross Margin", "85%", "88%", "90%"],
    ["Burn Multiple", "3.5x", "1.5x", "0.6x"],
]

col_widths = [Inches(3.3), Inches(2.8), Inches(2.8), Inches(2.8)]
table_left = Inches(1.1)
table_top = Inches(4.0)

for j, h in enumerate(headers):
    x = table_left + sum(w for w in col_widths[:j])
    add_tb(s, x, table_top, col_widths[j], Inches(0.3),
        h, size=10, bold=True, color=TEXT3, align=PP_ALIGN.CENTER)

for i, row in enumerate(rows):
    y = table_top + Inches(0.35 + i * 0.33)
    for j, cell in enumerate(row):
        x = table_left + sum(w for w in col_widths[:j])
        clr = TEXT if j == 0 else TEXT2
        bld = True if j == 0 else (True if 'Cr' in cell or 'L' in cell else False)
        add_tb(s, x, y, col_widths[j], Inches(0.3),
            cell, size=10, bold=bld, color=clr, align=PP_ALIGN.CENTER)

add_tb(s, Inches(1.1), Inches(6.85), Inches(11), Inches(0.2),
    "To reach 1,000 paid users: 0.05% penetration of SOM (2M) = 1,000 users → ₹1.2L MRR",
    size=11, bold=True, color=TEXT2, align=PP_ALIGN.CENTER)

# ========== SLIDE 6: HOW WE GET USERS ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
section_header(s, "GROWTH", "User Acquisition Strategy")

channels = [
    ("Product-Led Growth", "₹20/user", "25% of target", [
        "Free tier → viral sharing within teams",
        "Built-in referral: 'Share a note' exports with NoteMeet watermark",
        "WhatsApp integration: share summary → recipient clicks → downloads",
        "Meeting detection: app opens when you join a calendar event",
    ]),
    ("Founder-Led Sales", "₹50/user", "30% of target", [
        "Direct outreach to founder communities (IndieHackers, YC, 100x, SuperTeam)",
        "Twitter/X content: building in public, threads, product demos",
        "LinkedIn: DM startup founders with personalized demos",
        "Speaker at meetups: India AI, ProductHood, Founder meetups",
    ]),
    ("Paid Acquisition", "₹150/user", "15% of target", [
        "Google Ads: 'meeting notes AI', 'transcription tool India'",
        "Twitter/X ads targeted at Indian startup audience",
        "LinkedIn ads: targeting 'Engineering Manager' + 'India'",
        "ProductHunt launch + HackerNews",
    ]),
    ("Partnerships", "₹10/user", "30% of target", [
        "Zapier integration → discoverability in their ecosystem",
        "GCC partnerships: offer 50-100 free licenses to pilot",
        "Notion/Trello integration → notes land in their workflow",
        "Pre-installed on productivity-focused laptops (eventual)",
    ]),
]

for i, (channel, cac, pct, tactics) in enumerate(channels):
    col, row = i % 2, i // 2
    x = Inches(0.8 + col * 6.3)
    y = Inches(1.6 + row * 2.8)
    add_card(s, x, y, Inches(5.9), Inches(2.5))
    add_tb(s, x + Inches(0.3), y + Inches(0.2), Inches(3.5), Inches(0.3),
        channel, size=15, bold=True, color=TEXT)
    add_tb(s, x + Inches(3.8), y + Inches(0.2), Inches(1.8), Inches(0.3),
        f"CAC: {cac}", size=11, color=TEXT3, align=PP_ALIGN.RIGHT)
    add_tb(s, x + Inches(4.5), y + Inches(0.5), Inches(1.1), Inches(0.25),
        f"{pct}", size=10, color=TEXT2, align=PP_ALIGN.RIGHT)
    tf = add_tb(s, x + Inches(0.3), y + Inches(0.6), Inches(5.3), Inches(1.7), "", size=11, color=TEXT2)
    for t in tactics:
        add_bullet(tf, f"→  {t}", size=10, color=TEXT2)

# ========== SLIDE 7: TEAM BUILDING ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
section_header(s, "PEOPLE", "How We Build the Team")

# Phase 1
add_card(s, Inches(0.8), Inches(1.5), Inches(3.7), Inches(5.3))
add_tb(s, Inches(1.1), Inches(1.65), Inches(3.1), Inches(0.3),
    "Phase 1: Solo + Contractors", size=14, bold=True, color=TEXT)
add_tb(s, Inches(1.1), Inches(2.0), Inches(3.1), Inches(0.2),
    "Pre-revenue to ₹2L MRR  |  ₹0-5L burn", size=10, color=TEXT3)
tf = add_tb(s, Inches(1.1), Inches(2.3), Inches(3.1), Inches(4.3), "", size=11, color=TEXT2)
for item in [
    "You (full-stack) — product, code, design, PM",
    "Contract UI designer — ₹50k/month retainer",
    "Contract React dev — ₹80k/month (part-time)",
    "Contract QA — ₹30k/month (part-time)",
    "Freelance content writer — blog + Twitter",
    "Tools: Cursor, Figma, Vercel, GitHub, Linear",
    "Google Cloud credits (startup program)",
    "Target: ship mobile MVP in 3 months",
]:
    add_bullet(tf, f"→  {item}", size=10, color=TEXT2)

# Phase 2
add_card(s, Inches(4.8), Inches(1.5), Inches(3.7), Inches(5.3))
add_tb(s, Inches(5.1), Inches(1.65), Inches(3.1), Inches(0.3),
    "Phase 2: Core Team", size=14, bold=True, color=TEXT)
add_tb(s, Inches(5.1), Inches(2.0), Inches(3.1), Inches(0.2),
    "₹2L-₹10L MRR  |  ₹5-20L burn", size=10, color=TEXT3)
tf = add_tb(s, Inches(5.1), Inches(2.3), Inches(3.1), Inches(4.3), "", size=11, color=TEXT2)
for item in [
    "You + Co-founder (technical/ML preferred)",
    "Hire #1: Mobile dev (React Native / Flutter)",
    "Hire #2: Backend/ML engineer",
    "Hire #3: Growth marketer",
    "Equity: 10-15% pool for first hires",
    "Office: Remote-first, co-working in Bangalore",
    "Salaries: ₹8-15L/year (below market + equity)",
    "Target: feature parity across platforms",
]:
    add_bullet(tf, f"→  {item}", size=10, color=TEXT2)

# Phase 3
add_card(s, Inches(8.8), Inches(1.5), Inches(3.7), Inches(5.3))
add_tb(s, Inches(9.1), Inches(1.65), Inches(3.1), Inches(0.3),
    "Phase 3: Scale Team", size=14, bold=True, color=TEXT)
add_tb(s, Inches(9.1), Inches(2.0), Inches(3.1), Inches(0.2),
    "₹10L+ MRR  |  ₹20L+ burn", size=10, color=TEXT3)
tf = add_tb(s, Inches(9.1), Inches(2.3), Inches(3.1), Inches(4.3), "", size=11, color=TEXT2)
for item in [
    "Team of 12-15: 5 eng, 2 design, 3 growth, 2 sales, 2 ops",
    "VP Engineering — own platform scaling",
    "ML engineer — fine-tune Whisper for Hinglish",
    "Enterprise sales — GCC/enterprise contracts",
    "Customer success — onboarding + support",
    "Office: Bangalore (Koramangala/Indiranagar)",
    "Salaries: market rates + meaningful equity",
    "Target: 50K paid users, enterprise pilots",
]:
    add_bullet(tf, f"→  {item}", size=10, color=TEXT2)

# ========== SLIDE 8: TECHNICAL CHALLENGES ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
section_header(s, "RISKS — TECHNICAL", "The hard technical problems")

challenges = [
    ("01", "Mobile Call Recording (OS Blocked)",
     "CRITICAL — Apple and Google have deliberately blocked third-party call recording. iOS: no call audio API. Android: VOICE_CALL source removed in Android 10. There is no clean workaround. The industry's answer is 'desktop-only.'",
     "SIP-based in-app dialer (user calls through our app). This is how Truecaller does it. Needs telecom partnerships. For v1: focus on desktop + in-person mic recording.", RED),
    ("02", "On-Device Whisper Performance",
     "Whisper ggml-base runs well on M-series Macs (2-3s for 30s audio). Mid-range Android phones (Snapdragon 7 series) struggle (15-20s latency). iPhone 13+ processes in 5-8s.",
     "Strategy: Distil Whisper (faster, 95% accuracy). Offer cloud transcription as fallback. Start with higher-end devices and expand downmarket as hardware improves.", TEXT),
    ("03", "Hinglish / Code-Switching Accuracy",
     "Standard Whisper has ~60% WER on clean English, but 35-40% on Hinglish. Indian accents, mixed languages, and domain jargon (startup lingo) cause degradation.",
     "Fine-tune Whisper on Indian meeting data (1000+ hours). Use LoRA adapters for specific domains. Build dataset from real usage (with user permission).", TEXT),
    ("04", "Offline-First Sync Architecture",
     "Users expect to record on phone, view on laptop. CRDT-based sync is complex. Conflict resolution for note edits. Encrypted transfer between devices.",
     "Use local-first approach (Automerge/Y.js). End-to-end encryption. Server acts as relay, not store. iPad as primary mobile target (fewer OS restrictions).", TEXT),
    ("05", "System Audio Capture Fragility",
     "macOS screen recording permission changes. Windows WASAPI loopback has edge cases. Linux PipeWire/PulseAudio differences. Each OS update risks breakage.",
     "CI/CD tests on beta OS versions. Crash reporting (Sentry). Graceful degradation: fall back to mic if system audio unavailable.", TEXT),
]

for i, (num, title, problem, solution, accent) in enumerate(challenges):
    y = Inches(1.5 + i * 1.15)
    add_card(s, Inches(0.8), y, Inches(11.7), Inches(1.0))
    add_tb(s, Inches(1.0), y + Inches(0.08), Inches(0.3), Inches(0.25),
        num, size=11, bold=True, color=TEXT3)
    add_tb(s, Inches(1.3), y + Inches(0.05), Inches(4.5), Inches(0.25),
        title, size=13, bold=True, color=TEXT)
    add_tb(s, Inches(1.0), y + Inches(0.32), Inches(11.3), Inches(0.35),
        problem, size=9, color=TEXT2)
    add_tb(s, Inches(1.0), y + Inches(0.68), Inches(11.3), Inches(0.3),
        f"Strategy: {solution}", size=9, bold=False, color=TEXT)

# ========== SLIDE 9: REGULATORY CHALLENGES ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
section_header(s, "RISKS — REGULATORY & LEGAL", "The compliance landscape")

regs = [
    ("📜", "DPDP Act 2023 (India)", "HIGH",
     "Mandates data localization, consent management, and data protection impact assessments. Penalties up to ₹250 Cr.",
     "Our approach: Keep data on-device by default. No cloud storage. Consent-first UI. Indian servers when cloud sync is added (AWS Mumbai). Appoint Data Protection Officer at scale."),
    ("🎧", "Call Recording Consent Laws", "HIGH",
     "India: single-party consent (you can record your own calls). US: varies by state (some require two-party). EU: requires explicit consent (GDPR Art. 5).",
     "Our approach: Always show recording indicator. Require explicit tap to record. Never auto-record. For team plans: admin configures compliance rules per jurisdiction."),
    ("🔐", "Audio Data & Privacy", "MEDIUM",
     "Transcripts contain PII (names, contact info, financial details). On-device processing means raw audio never leaves. But cloud LLM calls send text summaries.",
     "Our approach: Local Whisper for transcription. Cloud LLM only with anonymized text. Opt-in for cloud features. Enterprise: completely on-premise option."),
    ("🏪", "App Store Policies", "MEDIUM",
     "Apple requires apps with call recording to prove consent. Google Play banned call recording via Accessibility Service (May 2023). Both review recording features strictly.",
     "Our approach: Desktop-first (no app store gatekeeping). Mobile app: mic recording only (in-person meetings). In-app dialer reviewed as 'communications app,' not 'recorder.'"),
    ("💳", "Payment Processing", "LOW",
     "Need PCI DSS compliance for card payments. Razorpay/Stripe handle this, but recurring billing requires proper webhook handling and invoice generation.",
     "Our approach: Razorpay for India (they handle GST compliance). Stripe for international. Automated invoicing. Self-serve + enterprise manual billing."),
]

for i, (icon, title, severity, problem, approach) in enumerate(regs):
    y = Inches(1.5 + i * 1.1)
    sev_color = RED if severity == "HIGH" else TEXT2
    add_card(s, Inches(0.8), y, Inches(11.7), Inches(0.95))
    add_tb(s, Inches(1.0), y + Inches(0.05), Inches(0.35), Inches(0.25),
        icon, size=16)
    add_tb(s, Inches(1.4), y + Inches(0.05), Inches(4.0), Inches(0.25),
        title, size=13, bold=True, color=TEXT)
    add_tb(s, Inches(5.4), y + Inches(0.05), Inches(1.0), Inches(0.25),
        severity, size=10, bold=True, color=sev_color)
    add_tb(s, Inches(1.0), y + Inches(0.35), Inches(5.0), Inches(0.25),
        f"Risk: {problem}", size=9, color=TEXT2)
    add_tb(s, Inches(1.0), y + Inches(0.63), Inches(5.0), Inches(0.25),
        f"Plan: {approach}", size=9, color=TEXT)

# ========== SLIDE 10: COMPETITIVE LANDSCAPE ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
section_header(s, "COMPETITION", "Who else is in this space?")

# Competitor cards
comps = [
    ("Granola", "$125M valuation", "Desktop AI notepad. No bot. £15/mo. iOS-only mobile. No Indian languages. No WhatsApp. Great UX but expensive and India-unfriendly."),
    ("Otter.ai", "$500M+ raised", "Bot joins meetings. $17/mo. English only. No WhatsApp. Good transcription but invasive and useless for Indian meetings."),
    ("Fireflies.ai", "$75M+ raised", "Bot joins meetings. $18/mo. Limited Indian language support. No WhatsApp. Strong enterprise but awkward bot approach."),
    ("Talat", "Early stage (India)", "Desktop AI notepad. Hindi transcription. No WhatsApp. No mobile. Free. Promising but very early and limited feature set."),
    ("Sutra AI", "Early stage (India)", "10 Indian languages. ₹399/60min pay-per-use. No WhatsApp. No mobile. Expensive for heavy users. Good language support."),
    ("Built-in (Apple/Google)", "Platform feature", "iOS 18.1+ has call recording. Google Recorder on Pixel. Limited to their ecosystems. No cross-platform. No AI note generation."),
]

for i, (name, funding, analysis) in enumerate(comps):
    col, row = i % 3, i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(1.5 + row * 2.8)
    add_card(s, x, y, Inches(3.7), Inches(2.5))
    add_tb(s, x + Inches(0.25), y + Inches(0.2), Inches(2.5), Inches(0.3),
        name, size=15, bold=True, color=TEXT)
    add_tb(s, x + Inches(2.75), y + Inches(0.2), Inches(0.7), Inches(0.3),
        funding, size=8, color=TEXT3, align=PP_ALIGN.RIGHT)
    add_tb(s, x + Inches(0.25), y + Inches(0.6), Inches(3.2), Inches(1.7),
        analysis, size=10, color=TEXT2)

# ========== SLIDE 11: FINANCIAL PROJECTIONS ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
section_header(s, "FINANCIALS", "Detailed Financial Projection")

# Revenue breakdown
add_card(s, Inches(0.8), Inches(1.5), Inches(5.7), Inches(5.5))
add_tb(s, Inches(1.1), Inches(1.65), Inches(5.1), Inches(0.3),
    "Revenue Breakdown (Year 3 target: ₹6Cr ARR)", size=13, bold=True, color=TEXT)

rev_items = [
    ("Pro Subscriptions", "₹4.2 Cr", "70%", "50K users × ₹199/mo × 12 × 0.4 (40% conversion)"),
    ("Business Subscriptions", "₹1.5 Cr", "25%", "2K teams × ₹599/user/mo × 12 × 3 users avg"),
    ("API / White-label", "₹0.3 Cr", "5%", "Enterprise licensing, custom deployments"),
]
tf = add_tb(s, Inches(1.1), Inches(2.1), Inches(5.1), Inches(1.5), "", size=10, color=TEXT2)
for item, amount, pct, note in rev_items:
    add_bullet(tf, f"{item}  ({pct})", size=11, bold=True, color=TEXT)
    add_bullet(tf, f"   {amount}  —  {note}", size=10, color=TEXT2)

add_tb(s, Inches(1.1), Inches(4.0), Inches(5.1), Inches(0.3),
    "Cost Structure (Year 3)", size=13, bold=True, color=TEXT)
costs = [
    ("Engineering (8 people)", "₹1.2 Cr", "20%"),
    ("Cloud Infrastructure", "₹0.5 Cr", "8%"),
    ("LLM API Costs", "₹0.8 Cr", "13%"),
    ("Marketing & Sales", "₹1.0 Cr", "17%"),
    ("Office & Operations", "₹0.5 Cr", "8%"),
    ("Legal & Compliance", "₹0.2 Cr", "3%"),
]
tf = add_tb(s, Inches(1.1), Inches(4.4), Inches(5.1), Inches(2.0), "", size=10, color=TEXT2)
for item, amount, pct in costs:
    add_bullet(tf, f"{item}  {amount}  ({pct})", size=10, color=TEXT2)

add_tb(s, Inches(1.1), Inches(6.3), Inches(5.1), Inches(0.5),
    "Gross Margin: 90%  |  Net Margin: 30%  |  Burn multiple: 0.6x",
    size=11, bold=True, color=TEXT)

# Funding needs
add_card(s, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.5))
add_tb(s, Inches(7.1), Inches(1.65), Inches(5.1), Inches(0.3),
    "Funding Requirements", size=13, bold=True, color=TEXT)

rounds = [
    ("Pre-Seed (Now)", "₹25-50L ($30-60k)", [
        "12-18 months runway",
        "Mobile app development (iOS + Android)",
        "Cloud infra + API credits",
        "Legal + compliance setup",
        "Part-time contractors",
    ]),
    ("Seed (Year 1 end)", "₹1.5-2.5 Cr ($180-300k)", [
        "Core team of 4-5",
        "Full-time founder salaries",
        "Marketing push (paid acquisition)",
        "Fine-tune Whisper for Indian languages",
        "Enterprise pilot program",
    ]),
    ("Series A (Year 2-3)", "₹8-12 Cr ($1-1.5M)", [
        "Scale team to 12-15",
        "Office in Bangalore",
        "Enterprise sales team",
        "On-premise deployment capability",
        "Data center in India",
    ]),
]

for i, (round_name, amount, uses) in enumerate(rounds):
    y = Inches(1.95 + i * 1.7)
    add_tb(s, Inches(7.1), y, Inches(3.5), Inches(0.25),
        round_name, size=12, bold=True, color=TEXT)
    add_tb(s, Inches(7.1), y + Inches(0.25), Inches(3.5), Inches(0.2),
        amount, size=10, color=TEXT2)
    tf = add_tb(s, Inches(7.1), y + Inches(0.5), Inches(5.1), Inches(1.0), "", size=9, color=TEXT2)
    for u in uses:
        add_bullet(tf, f"→  {u}", size=9, color=TEXT2)

add_tb(s, Inches(7.1), Inches(6.6), Inches(5.1), Inches(0.3),
    "Target investors: Indian angels (better understanding of market) + AI-focused deep tech funds",
    size=10, bold=False, color=TEXT2)

# ========== SLIDE 12: WHY YOU (CO-FOUNDER FIT) ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
section_header(s, "THE ASK", "Why you? Why now?")

add_card(s, Inches(0.8), Inches(1.5), Inches(5.7), Inches(5.0))
add_tb(s, Inches(1.1), Inches(1.7), Inches(5.1), Inches(0.3),
    "What I'm Looking For", size=16, bold=True, color=TEXT)
tf = add_tb(s, Inches(1.1), Inches(2.2), Inches(5.1), Inches(4.0), "", size=12, color=TEXT2)
for item in [
    "A technical co-founder (ideally ML/backend)",
    "OR a business co-founder (sales/marketing/growth)",
    "Someone who lives and breathes the Indian startup ecosystem",
    "Willing to start part-time, go full-time when revenue hits ₹2L MRR",
    "Comfortable with high risk, low pay, high equity",
    "Based in Bangalore (preferred, not required)",
    "Equity: 20-40% depending on role and commitment",
    "Salary: ₹0 until funded or revenue covers it",
    "Timeline: want to decide in the next 2 weeks",
]:
    add_bullet(tf, f"→  {item}", size=12, color=TEXT2)

add_card(s, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.0))
add_tb(s, Inches(7.1), Inches(1.7), Inches(5.1), Inches(0.3),
    "Why Join Now?", size=16, bold=True, color=TEXT)
tf = add_tb(s, Inches(7.1), Inches(2.2), Inches(5.1), Inches(4.0), "", size=12, color=TEXT2)
for item in [
    "v1 already shipping — real users, real feedback",
    "Only ₹59k spent to reach MVP (capital efficient)",
    "Category validated: Granola $125M, Otter $500M+ raised",
    "India is completely underserved: no competitor does WhatsApp + Indian languages + mobile",
    "DPDP Act creates urgency for India-first solutions",
    "WhatsApp Business has 200M+ Indian users — massive unmet need",
    "On-device AI is finally feasible (Whisper, Llama 3)",
    "First-mover advantage in 'AI notepad for Indian teams'",
    "We can bootstrap to ₹2L MRR before needing external funding",
]:
    add_bullet(tf, f"→  {item}", size=12, color=TEXT2)

# ========== SLIDE 13: NEXT STEPS ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

add_tb(s, Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.0),
    "What happens next?", size=42, bold=True, color=TEXT, align=PP_ALIGN.CENTER)

steps = [
    ("1", "Review this deck", "Take a week. Think about it. Ask hard questions."),
    ("2", "Try the app", "I'll send you a build. Use it for your own meetings. See if it clicks."),
    ("3", "Honest conversation", "We talk about what you liked, what you hated, what scares you."),
    ("4", "Decide together", "If yes: we agree on terms, equity, and start. If no: no hard feelings."),
    ("5", "Ship mobile MVP", "Target: Android + iOS in 3 months."),
    ("6", "Launch at ₹199/mo", "Start charging from month 4. Validate willingness to pay."),
]

for i, (num, title, desc) in enumerate(steps):
    y = Inches(3.2 + i * 0.65)
    add_tb(s, Inches(2.0), y, Inches(0.4), Inches(0.4),
        num, size=18, bold=True, color=TEXT2, align=PP_ALIGN.CENTER)
    add_tb(s, Inches(2.6), y, Inches(3.5), Inches(0.4),
        title, size=16, bold=True, color=TEXT)
    add_tb(s, Inches(6.2), y, Inches(5.5), Inches(0.4),
        desc, size=13, color=TEXT2)

add_card(s, Inches(2), Inches(6.5), Inches(9.3), Inches(0.65), SURFACE3, BORDER2)
add_tb(s, Inches(2.3), Inches(6.58), Inches(8.7), Inches(0.4),
    "Built solo in 8 weeks.  ₹59k cost.  52 tests passing.  Now I need a partner to take it further.",
    size=13, color=TEXT2, align=PP_ALIGN.CENTER)

# Save
output = os.path.expanduser("~/Desktop/NoteMeet_CoFounder_Deck.pptx")
prs.save(output)
print(f"✅ Saved to {output}")
print(f"   Slides: {len(prs.slides)}")
